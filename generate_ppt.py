from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# Helper for title slide
def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]
    title_box.text = title
    subtitle_box.text = subtitle
    return slide

# Helper for bullet slide
def add_bullet_slide(title, points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    title_box.text = title
    
    tf = slide.placeholders[1].text_frame
    tf.text = points[0]
    for point in points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
    return slide

# Slide 1: Welcome
add_title_slide("Welcome to RoadWatch AI", "Public Audit System & Infrastructure Transparency\nTeam Submission")

# Slide 2: Problem Statement
add_bullet_slide("Problem Statement", [
    "Lack of transparency in municipal road repair budgets.",
    "Citizens have no reliable way to report or track road hazards.",
    "Road conditions deteriorate rapidly without timely intervention.",
    "No gamified incentive for citizens to contribute to road safety."
])

# Slide 3: Our Solution
add_bullet_slide("Our Solution: RoadWatch AI", [
    "A real-time anomaly detection platform powered by dashcam feeds.",
    "Dynamic mapping of potholes, cracks, and faded markings.",
    "Direct synchronization with municipal budgets to highlight expenditure efficiency.",
    "Citizen rewards program to incentivize continuous public auditing."
])

# Slide 4: Key Features & Demo
add_bullet_slide("Key Features", [
    "Live Map: Glowing severity-based markers for instantaneous hazard identification.",
    "Citizen Rewards: Experience points and rank badges for verified anomaly reporting.",
    "Road Analytics: Comprehensive indices including Ride Comfort and Surface Quality.",
    "Full Anomaly Reports: Actionable telemetry for rapid repair crew dispatch."
])

# Slide 5: Technical Stack
add_bullet_slide("Technical Stack", [
    "Frontend: React.js, Vite, Tailwind CSS (Glassmorphism UI).",
    "Mapping: React-Leaflet, CartoDB Dark Base Maps.",
    "Backend & AI: Python FastAPI (Simulated Dashcam Inference).",
    "State Management: React Hooks and Axios."
])

# Slide 6: Social & Financial Impact
add_bullet_slide("Impact & Future Scope", [
    "Increases fiscal transparency by correlating road budgets with actual quality.",
    "Reduces vehicular damage and accidents through early hazard detection.",
    "Encourages crowdsourced governance via the Citizen Rewards system.",
    "Future: Integration with actual vehicle telematics and blockchain ledgers."
])

# Slide 7: Thank You
add_title_slide("Thank You!", "Q&A\n\n(Code & Database submitted via Unstop)")

prs.save("RoadWatch_AI_Submission.pptx")
print("Presentation saved as RoadWatch_AI_Submission.pptx")
